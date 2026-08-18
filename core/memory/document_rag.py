"""RAG по пользовательским файлам (PDF/текст) на базе ChromaDB.

Сканирует ``settings.paths.documents_dir``, режет текст на чанки и
индексирует в отдельную коллекцию ``documents`` (эмбеддинги — общий
проектный :class:`Embedder`, как и у долгой памяти).

Обработка ошибок: битый PDF, неподдерживаемый формат или пустой файл —
логируются и пропускаются, не прерывая индексацию остальных файлов.
"""

from __future__ import annotations

import hashlib
from pathlib import Path
from typing import List, Optional, Set

from config.settings import Settings
from core.memory.embedder import Embedder
from core.utils.logger import get_logger

__all__ = ["DocumentRAG", "read_pdf", "read_text_file", "read_document", "chunk_text", "DOCUMENTS_COLLECTION"]

log = get_logger(__name__)

#: Поддерживаемые расширения для индексации.
_SUPPORTED_EXTENSIONS = (".txt", ".md", ".pdf", ".text", ".docx", ".xlsx", ".pptx", ".png", ".jpg", ".jpeg", ".webp")

#: Имя коллекции ChromaDB для документов.
DOCUMENTS_COLLECTION = "documents"


def read_pdf(path: Path) -> str:
    """Извлекает текст из PDF через pypdf.

    Raises:
        RuntimeError: если файл не читается как PDF или pypdf недоступен.
    """
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise RuntimeError("pypdf не установлен — чтение PDF недоступно") from exc

    try:
        reader = PdfReader(str(path))
    except Exception as exc:
        raise RuntimeError(f"Не удалось открыть PDF {path}: {exc}") from exc

    pages: List[str] = []
    for page in reader.pages:
        try:
            text = page.extract_text() or ""
        except Exception as exc:
            log.debug("Пропущена страница в %s: %s", path, exc)
            text = ""
        if text.strip():
            pages.append(text)
    return "\n\n".join(pages)


def read_text_file(path: Path) -> str:
    """Читает текстовый файл (.txt/.md/.text). Декодирует UTF-8, fallback на cp1251."""
    try:
        return path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        return path.read_text(encoding="cp1251", errors="replace")
    except OSError as exc:
        raise RuntimeError(f"Не удалось прочитать файл {path}: {exc}") from exc


def read_document(path: Path) -> str:
    """Extract text from local text, office and image formats."""
    suffix = path.suffix.casefold()
    if suffix in {".txt", ".md", ".text"}:
        return read_text_file(path)
    if suffix == ".pdf":
        return read_pdf(path)
    if suffix == ".docx":
        try:
            from docx import Document
        except ImportError as exc:
            raise RuntimeError("python-docx не установлен") from exc
        return "\n".join(p.text for p in Document(str(path)).paragraphs if p.text.strip())
    if suffix == ".xlsx":
        try:
            from openpyxl import load_workbook
        except ImportError as exc:
            raise RuntimeError("openpyxl не установлен") from exc
        book = load_workbook(str(path), read_only=True, data_only=True)
        rows: list[str] = []
        for sheet in book.worksheets:
            rows.append(f"[sheet: {sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                values = [str(item) for item in row if item is not None]
                if values:
                    rows.append(" | ".join(values))
        return "\n".join(rows)
    if suffix == ".pptx":
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise RuntimeError("python-pptx не установлен") from exc
        slides: list[str] = []
        for index, slide in enumerate(Presentation(str(path)).slides, start=1):
            texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            if texts:
                slides.append(f"[slide: {index}]\n" + "\n".join(texts))
        return "\n\n".join(slides)
    if suffix in {".png", ".jpg", ".jpeg", ".webp"}:
        try:
            import pytesseract
            from PIL import Image
        except ImportError as exc:
            raise RuntimeError("pytesseract/Pillow не установлены") from exc
        return str(pytesseract.image_to_string(Image.open(path), lang="rus+eng") or "")
    raise RuntimeError(f"Формат не поддерживается: {path.suffix}")


def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> List[str]:
    """Режет текст на перекрывающиеся чанки по словам.

    Args:
        text: исходный текст.
        chunk_size: целевой размер чанка в символах.
        overlap: перекрытие между соседними чанками (символов).

    Returns:
        Список непустых чанков. Пустой при пустом входе.
    """
    if not text or not text.strip():
        return []
    if chunk_size <= 0:
        chunk_size = 500
    if overlap < 0:
        overlap = 0
    if overlap >= chunk_size:
        overlap = chunk_size // 4

    words = text.split()
    if not words:
        return []

    chunks: List[str] = []
    current = ""
    for word in words:
        if current and len(current) + len(word) + 1 > chunk_size:
            chunks.append(current.strip())
            # Начинаем новый чанк с перекрытием из конца предыдущего.
            overlap_text = current[-overlap:] if overlap else ""
            current = (overlap_text + " " + word).strip()
        else:
            current = (current + " " + word).strip() if current else word
    if current.strip():
        chunks.append(current.strip())
    return chunks


class DocumentRAG:
    """Индексатор и поиск по документам пользователя."""

    def __init__(self, settings: Settings, embedder: Embedder) -> None:
        """
        Args:
            settings: конфигурация (берётся ``documents_dir``).
            embedder: общий эмбеддер проекта.
        """
        self._settings = settings
        self._embedder = embedder
        self._client = None
        self._collection = None
        self._initialized = False
        # Запоминаем уже проиндексированные файлы, чтобы index_all не дублировал.
        self._indexed: Set[str] = set()
        self._indexed_mtimes: dict[str, int] = {}
        self._init_client()

    # ------------------------------------------------------------------ #
    #  ChromaDB
    # ------------------------------------------------------------------ #

    def _init_client(self) -> None:
        try:
            import chromadb
        except ImportError as exc:
            log.error("chromadb не установлен: RAG-документы недоступны: %s", exc)
            return

        # Подавляем сломанную телеметрию ChromaDB (см. long_term.py).
        from core.memory.long_term import disable_chroma_telemetry_noise
        disable_chroma_telemetry_noise()

        try:
            documents_dir = self._settings.paths.resolved("documents_dir")
            documents_dir.mkdir(parents=True, exist_ok=True)
            # anonymized_telemetry=False — отключаем сетевую телеметрию
            # ChromaDB (шум в логах), функциональность RAG не меняется.
            self._client = chromadb.PersistentClient(
                path=str(documents_dir / ".chroma"),
                settings=chromadb.Settings(anonymized_telemetry=False),
            )
            self._collection = self._client.get_or_create_collection(
                name=DOCUMENTS_COLLECTION,
                metadata={"hnsw:space": "cosine"},
            )
            self._initialized = True
            log.info("RAG-документы инициализированы: %s", documents_dir)
        except Exception as exc:
            log.error("Не удалось инициализировать RAG ChromaDB: %s", exc)
            self._initialized = False

    def _ensure(self) -> bool:
        if self._initialized and self._collection is not None:
            return True
        if self._client is None:
            self._init_client()
        return self._initialized and self._collection is not None

    # ------------------------------------------------------------------ #
    #  Индексация
    # ------------------------------------------------------------------ #

    def index_file(self, path: Path) -> int:
        """Индексирует один файл. Возвращает число добавленных чанков.

        Args:
            path: путь к .txt/.md/.pdf файлу.

        Returns:
            Количество успешно добавленных чанков (0 при ошибке/пустом файле).
        """
        path = Path(path)
        if not self._ensure():
            return 0
        if not path.is_file():
            log.warning("Файл не существует: %s", path)
            return 0
        if path.suffix.lower() not in _SUPPORTED_EXTENSIONS:
            log.debug("Пропущен неподдерживаемый формат: %s", path)
            return 0

        try:
            text = read_document(path)
        except Exception as exc:
            log.warning("Не удалось прочитать %s: %s — пропускаю", path, exc)
            return 0

        chunks = chunk_text(text)
        if not chunks:
            log.debug("Файл пуст после извлечения текста: %s", path)
            return 0

        # Стабильный id чанка, чтобы повторная индексация не дублировала.
        file_key = str(path.resolve())
        try:
            ids = [f"{hashlib.md5(file_key.encode()).hexdigest()}-{i}" for i in range(len(chunks))]
            digest = hashlib.sha256(path.read_bytes()).hexdigest()
            metadatas = [{
                "source": file_key,
                "chunk_index": i,
                "sha256": digest,
                "mtime_ns": path.stat().st_mtime_ns,
                "format": path.suffix.casefold().lstrip("."),
            } for i in range(len(chunks))]
            if file_key in self._indexed:
                try:
                    self._collection.delete(ids=ids)
                except Exception:
                    pass
            self._collection.add(documents=chunks, ids=ids, metadatas=metadatas)
            self._indexed.add(file_key)
            self._indexed_mtimes[file_key] = path.stat().st_mtime_ns
            log.info("Проиндексирован %s: %d чанков", path.name, len(chunks))
            return len(chunks)
        except Exception as exc:
            log.error("Ошибка индексации %s: %s", path, exc)
            return 0

    def index_all(self) -> dict:
        """Сканирует documents_dir и индексирует файлы, ещё не проиндексированные.

        Returns:
            Словарь ``{"indexed": int, "chunks": int, "skipped": int}``.
        """
        if not self._ensure():
            return {"indexed": 0, "chunks": 0, "skipped": 0}

        documents_dir = self._settings.paths.resolved("documents_dir")
        files = [
            p for p in documents_dir.iterdir()
            if p.is_file() and p.suffix.lower() in _SUPPORTED_EXTENSIONS
        ]

        indexed_count = 0
        chunks_total = 0
        skipped = 0
        for file_path in files:
            file_key = str(file_path.resolve())
            if (file_key in self._indexed and
                    self._indexed_mtimes.get(file_key) == file_path.stat().st_mtime_ns):
                skipped += 1
                continue
            added = self.index_file(file_path)
            if added > 0:
                indexed_count += 1
                chunks_total += added
            else:
                skipped += 1

        log.info(
            "index_all: проиндексировано файлов=%d, чанков=%d, пропущено=%d",
            indexed_count, chunks_total, skipped,
        )
        return {"indexed": indexed_count, "chunks": chunks_total, "skipped": skipped}

    # ------------------------------------------------------------------ #
    #  Поиск
    # ------------------------------------------------------------------ #

    def search_documents(self, query: str, top_k: int = 3) -> List[str]:
        """Семантический поиск по проиндексированным документам.

        Args:
            query: поисковый запрос.
            top_k: сколько фрагментов вернуть.

        Returns:
            Список найденных текстовых фрагментов (пустой при отсутствии/ошибке).
        """
        if not query or not query.strip():
            return []
        if not self._ensure():
            return []

        try:
            results = self._collection.query(
                query_texts=[query],
                n_results=max(1, top_k),
            )
        except Exception as exc:
            log.error("Ошибка поиска по документам: %s", exc)
            return []

        documents = (results or {}).get("documents") or [[]]
        # §22 — RAG-фрагменты это внешние ДАННЫЕ. Единая точка обёртки:
        # оборачиваем каждый чанк на границе поиск→модель (wrap_untrusted
        # идемпотентен — если вызывается повторно, дубля не будет).
        from core.safety import wrap_untrusted
        return [wrap_untrusted(str(doc), source="RAG") for doc in documents[0] if doc]
